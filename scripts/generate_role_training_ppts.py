from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

import generate_training_ppt as base


ROOT = Path(r"C:\Users\baek_\newmso")
MANIFEST_DIR = ROOT / "docs" / "manuals" / "ppt_assets" / "manifest"
OUTPUT_DIR = ROOT / "docs" / "manuals" / "generated"


ROLE_DECKS = {
    "employee": {
        "title": "NEWMSO 사원용 화면 설명서",
        "subtitle": "직원이 실제로 가장 자주 쓰는 기능을 화면 중심으로 설명합니다.",
        "output": OUTPUT_DIR / "NEWMSO_사원용_화면설명서.pptx",
        "slides": ["01_login", "02_main_shell", "03_mypage", "04_chat", "05_board", "06_approval"],
        "summary": [
            "핵심 사용 메뉴: 내정보, 채팅, 게시판, 전자결재",
            "가장 자주 하는 업무: 출퇴근 확인, 증명서 신청, 서류 제출, 연차 신청",
            "메뉴가 안 보이면 권한 문제일 수 있으므로 부서장 또는 관리자에게 문의",
        ],
    },
    "manager": {
        "title": "NEWMSO 부서장용 화면 설명서",
        "subtitle": "부서장이 승인과 팀 운영에 집중할 수 있도록 핵심 화면을 모았습니다.",
        "output": OUTPUT_DIR / "NEWMSO_부서장용_화면설명서.pptx",
        "slides": ["01_login", "02_main_shell", "10_mypage_manager", "11_approval_manager", "07_hr", "08_inventory"],
        "summary": [
            "핵심 사용 메뉴: 내정보, 전자결재, 인사관리, 재고관리",
            "가장 자주 하는 업무: 승인 처리, 팀 현황 점검, 부족 재고 대응",
            "승인 결과는 실제 근태, 휴가, 구매 흐름에 영향을 주므로 주의",
        ],
    },
    "admin": {
        "title": "NEWMSO 관리자용 화면 설명서",
        "subtitle": "운영 담당자가 자주 쓰는 관리자 기능과 연계 화면을 모았습니다.",
        "output": OUTPUT_DIR / "NEWMSO_관리자용_화면설명서.pptx",
        "slides": ["01_login", "02_main_shell", "09_admin", "12_admin_permissions", "07_hr", "08_inventory", "06_approval"],
        "summary": [
            "핵심 사용 메뉴: 관리자, 인사관리, 재고관리, 전자결재",
            "중요 화면: 직원 권한, 백업/초기화, 팝업관리, 양식빌더",
            "권한 변경과 데이터 작업은 운영 영향이 크므로 변경 전후 확인 필요",
        ],
    },
    "master": {
        "title": "NEWMSO 마스터용 화면 설명서",
        "subtitle": "최고권한 운영자가 통제해야 하는 실제 화면과 운영 기준을 함께 정리했습니다.",
        "output": OUTPUT_DIR / "NEWMSO_마스터용_화면설명서.pptx",
        "slides": ["01_login", "02_main_shell", "09_admin", "12_admin_permissions"],
        "summary": [
            "현재 원본 구조상 마스터 전용 별도 UI는 없고, 관리자 화면을 최고 권한 범위로 운영",
            "핵심 통제 대상: 권한, 감사로그, 백업, 초기화, 양식/문서 구조 변경",
            "운영 정책과 고위험 작업 승인 기준 문서로 함께 사용 권장",
        ],
        "governance": [
            {
                "title": "마스터 운영 원칙",
                "bullets": [
                    "최소 권한 원칙: 필요한 메뉴만 열기",
                    "변경 기록 원칙: 권한/데이터/양식 변경 사유 기록",
                    "작업 전 백업 원칙: 구조 변경 전 백업 확보",
                    "운영/테스트 분리 원칙: 실험성 변경은 운영에서 바로 하지 않기",
                ],
            },
            {
                "title": "마스터 점검 체크리스트",
                "bullets": [
                    "일일: 장애 제보, 권한 사고, 결재 적체 여부 확인",
                    "주간: 감사로그, 접근감사로그, 백업 상태 점검",
                    "월간: 관리자 계정 목록, 오래된 권한, 핵심 데이터 이상치 검토",
                    "고위험 작업: 데이터초기화, 대량 권한 변경, 양식 구조 변경은 별도 승인 후 진행",
                ],
            },
        ],
    },
}


def load_manifests():
    items = {}
    for path in MANIFEST_DIR.glob("*.json"):
        data = json.loads(path.read_text(encoding="utf-8"))
        items[data["id"]] = data
    return items


def add_role_title_slide(prs: Presentation, title: str, subtitle: str, summary: list[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(245, 247, 250)

    tag = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.6),
        Inches(0.55),
        Inches(1.9),
        Inches(0.42),
    )
    tag.fill.solid()
    tag.fill.fore_color.rgb = RGBColor(37, 99, 235)
    tag.line.color.rgb = RGBColor(37, 99, 235)
    base.add_textbox(slide, Inches(0.8), Inches(0.61), Inches(1.5), Inches(0.2), "역할별 배포본", 15, True, (255, 255, 255))

    base.add_textbox(slide, Inches(0.62), Inches(1.35), Inches(7.2), Inches(0.7), title, 28, True, (15, 23, 42))
    base.add_textbox(slide, Inches(0.62), Inches(2.05), Inches(7.0), Inches(0.6), subtitle, 17, False, (71, 85, 105))

    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.62),
        Inches(3.0),
        Inches(5.8),
        Inches(2.8),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(255, 255, 255)
    panel.line.color.rgb = RGBColor(226, 232, 240)
    base.add_textbox(slide, Inches(0.9), Inches(3.28), Inches(3.4), Inches(0.35), "이 덱에서 중점적으로 보는 내용", 19, True, (30, 41, 59))
    base.add_bullets(slide, Inches(0.9), Inches(3.72), Inches(5.0), Inches(1.8), summary)

    right = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(6.8),
        Inches(1.2),
        Inches(5.8),
        Inches(5.5),
    )
    right.fill.solid()
    right.fill.fore_color.rgb = RGBColor(230, 238, 255)
    right.line.color.rgb = RGBColor(191, 219, 254)
    base.add_textbox(slide, Inches(7.1), Inches(1.55), Inches(3.2), Inches(0.4), "활용 방식", 20, True, (29, 78, 216))
    base.add_bullets(
        slide,
        Inches(7.1),
        Inches(2.0),
        Inches(4.9),
        Inches(3.4),
        [
            "신규 사용자 교육",
            "사내 공지 첨부 자료",
            "역할별 업무 설명",
            "권한 문의 응대 기준 자료",
        ],
    )


def add_text_slide(prs: Presentation, title: str, bullets: list[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(248, 250, 252)

    header = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.75),
        Inches(0.8),
        Inches(11.8),
        Inches(0.9),
    )
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(37, 99, 235)
    header.line.color.rgb = RGBColor(37, 99, 235)
    base.add_textbox(slide, Inches(1.05), Inches(1.05), Inches(8.2), Inches(0.3), title, 24, True, (255, 255, 255))

    body = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.9),
        Inches(2.0),
        Inches(11.5),
        Inches(4.8),
    )
    body.fill.solid()
    body.fill.fore_color.rgb = RGBColor(255, 255, 255)
    body.line.color.rgb = RGBColor(226, 232, 240)
    base.add_bullets(slide, Inches(1.2), Inches(2.35), Inches(10.7), Inches(4.0), bullets)


def build_role_deck(name: str, config: dict, manifests: dict):
    prs = Presentation()
    prs.slide_width = Inches(base.SLIDE_W)
    prs.slide_height = Inches(base.SLIDE_H)

    add_role_title_slide(prs, config["title"], config["subtitle"], config["summary"])

    for slide_id in config["slides"]:
        item = manifests[slide_id]
        annotated = base.annotate_image(item)
        base.add_content_slide(prs, item, annotated)

    for extra in config.get("governance", []):
        add_text_slide(prs, extra["title"], extra["bullets"])

    prs.save(str(config["output"]))


def main():
    base.ensure_dirs()
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    manifests = load_manifests()
    for name, config in ROLE_DECKS.items():
        build_role_deck(name, config, manifests)
        print(config["output"])


if __name__ == "__main__":
    main()
