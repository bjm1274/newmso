from __future__ import annotations

import json
import math
from pathlib import Path
from typing import Any

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(r"C:\Users\baek_\newmso")
ASSET_ROOT = ROOT / "docs" / "manuals" / "ppt_assets"
RAW_DIR = ASSET_ROOT / "raw"
MANIFEST_DIR = ASSET_ROOT / "manifest"
ANNOTATED_DIR = ASSET_ROOT / "annotated"
OUTPUT_DIR = ROOT / "docs" / "manuals" / "generated"
OUTPUT_PPTX = OUTPUT_DIR / "NEWMSO_원본기준_화면설명서_1차.pptx"


SLIDE_W = 13.333
SLIDE_H = 7.5


def ensure_dirs() -> None:
    ANNOTATED_DIR.mkdir(parents=True, exist_ok=True)
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)


def load_font(size: int, bold: bool = False) -> ImageFont.ImageFont:
    candidates = []
    if bold:
        candidates.extend(
            [
                r"C:\Windows\Fonts\malgunbd.ttf",
                r"C:\Windows\Fonts\맑은 고딕 Bold.ttf",
            ]
        )
    candidates.extend(
        [
            r"C:\Windows\Fonts\malgun.ttf",
            r"C:\Windows\Fonts\맑은 고딕.ttf",
            r"C:\Windows\Fonts\arial.ttf",
        ]
    )
    for candidate in candidates:
        if Path(candidate).exists():
            return ImageFont.truetype(candidate, size=size)
    return ImageFont.load_default()


def draw_label(draw: ImageDraw.ImageDraw, x: float, y: float, index: int, fill: tuple[int, int, int]) -> None:
    font = load_font(30, bold=True)
    radius = 22
    draw.ellipse((x - radius, y - radius, x + radius, y + radius), fill=fill, outline=(255, 255, 255), width=3)
    text = str(index)
    bbox = draw.textbbox((0, 0), text, font=font)
    tw = bbox[2] - bbox[0]
    th = bbox[3] - bbox[1]
    draw.text((x - tw / 2, y - th / 2 - 1), text, fill=(255, 255, 255), font=font)


def annotate_image(item: dict[str, Any]) -> Path:
    image_path = Path(item["image"])
    annotated_path = ANNOTATED_DIR / f"{item['id']}.png"

    image = Image.open(image_path).convert("RGBA")
    overlay = Image.new("RGBA", image.size, (255, 255, 255, 0))
    draw = ImageDraw.Draw(overlay)
    label_font = load_font(24, bold=True)
    note_font = load_font(18)

    palette = [
        (37, 99, 235),
        (22, 163, 74),
        (234, 88, 12),
        (168, 85, 247),
        (220, 38, 38),
        (14, 165, 233),
        (217, 119, 6),
        (236, 72, 153),
    ]

    for idx, hotspot in enumerate(item.get("hotspots", [])):
        color = palette[idx % len(palette)]
        box = hotspot["box"]
        x = box["x"]
        y = box["y"]
        w = box["width"]
        h = box["height"]

        draw.rounded_rectangle(
            (x, y, x + w, y + h),
            radius=16,
            outline=color + (255,),
            width=5,
            fill=color + (30,),
        )

        label_x = max(30, x - 16)
        label_y = max(30, y - 18)
        draw_label(draw, label_x, label_y, hotspot["index"], color)

        chip_text = hotspot["label"]
        chip_bbox = draw.textbbox((0, 0), chip_text, font=label_font)
        chip_w = (chip_bbox[2] - chip_bbox[0]) + 24
        chip_h = (chip_bbox[3] - chip_bbox[1]) + 14
        chip_left = min(max(12, x + 16), image.size[0] - chip_w - 12)
        chip_top = min(max(12, y + 12), image.size[1] - chip_h - 12)
        draw.rounded_rectangle(
            (chip_left, chip_top, chip_left + chip_w, chip_top + chip_h),
            radius=12,
            fill=(255, 255, 255, 235),
            outline=color + (255,),
            width=2,
        )
        draw.text((chip_left + 12, chip_top + 5), chip_text, fill=color + (255,), font=label_font)

        note_text = hotspot.get("note", "")
        if note_text:
            note_bbox = draw.textbbox((0, 0), note_text, font=note_font)
            note_w = (note_bbox[2] - note_bbox[0]) + 18
            note_h = (note_bbox[3] - note_bbox[1]) + 12
            note_left = min(max(12, x + 16), image.size[0] - note_w - 12)
            note_top = min(max(12, chip_top + chip_h + 8), image.size[1] - note_h - 12)
            draw.rounded_rectangle(
                (note_left, note_top, note_left + note_w, note_top + note_h),
                radius=10,
                fill=(255, 255, 255, 220),
                outline=(180, 180, 180, 255),
                width=1,
            )
            draw.text((note_left + 9, note_top + 4), note_text, fill=(55, 65, 81, 255), font=note_font)

    annotated = Image.alpha_composite(image, overlay).convert("RGB")
    annotated.save(annotated_path)
    return annotated_path


def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=(17, 24, 39), align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.name = "Malgun Gothic"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    return box


def add_bullets(slide, left, top, width, height, items):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    for idx, item in enumerate(items):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(16)
        p.font.name = "Malgun Gothic"
        p.font.color.rgb = RGBColor(55, 65, 81)
    return box


def add_title_slide(prs: Presentation, slide_count: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(245, 247, 250)

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.55),
        Inches(0.55),
        Inches(1.55),
        Inches(0.4),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(37, 99, 235)
    accent.line.color.rgb = RGBColor(37, 99, 235)
    add_textbox(slide, Inches(0.72), Inches(0.60), Inches(1.25), Inches(0.25), "원본 기준", 16, True, (255, 255, 255))

    add_textbox(
        slide,
        Inches(0.7),
        Inches(1.35),
        Inches(6.9),
        Inches(1.2),
        "NEWMSO 화면형 교육자료",
        28,
        True,
        (15, 23, 42),
    )
    add_textbox(
        slide,
        Inches(0.7),
        Inches(2.25),
        Inches(6.6),
        Inches(0.9),
        "실제 화면 캡처와 버튼 위치를 기준으로 정리한 1차 PPT입니다.",
        17,
        False,
        (71, 85, 105),
    )

    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.7),
        Inches(3.2),
        Inches(5.0),
        Inches(2.9),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(255, 255, 255)
    panel.line.color.rgb = RGBColor(226, 232, 240)
    add_bullets(
        slide,
        Inches(0.95),
        Inches(3.5),
        Inches(4.5),
        Inches(2.2),
        [
            "기준: C:\\Users\\baek_\\newmso 원본 프로그램",
            "형식: 실제 화면 캡처 + 번호 콜아웃",
            f"구성: 총 {slide_count + 1}장",
            "용도: 직원 교육, 부서장 안내, 운영자 브리핑",
        ],
    )

    right_panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(6.15),
        Inches(1.15),
        Inches(6.1),
        Inches(5.7),
    )
    right_panel.fill.solid()
    right_panel.fill.fore_color.rgb = RGBColor(230, 238, 255)
    right_panel.line.color.rgb = RGBColor(191, 219, 254)

    add_textbox(slide, Inches(6.45), Inches(1.55), Inches(5.3), Inches(0.5), "이번 1차 PPT 범위", 20, True, (29, 78, 216))
    add_bullets(
        slide,
        Inches(6.45),
        Inches(2.0),
        Inches(5.2),
        Inches(3.0),
        [
            "로그인",
            "메인 사이드바",
            "내정보",
            "채팅",
            "게시판",
            "전자결재",
            "인사관리",
            "재고관리",
            "관리자 화면",
        ],
    )
    add_textbox(slide, Inches(6.45), Inches(5.6), Inches(5.1), Inches(0.8), "역할별 세부 PPT는 이 1차 자료를 바탕으로 계속 확장할 수 있습니다.", 15, False, (55, 65, 81))


def fit_image(slide, image_path: Path, left: float, top: float, width: float, height: float):
    with Image.open(image_path) as image:
        image_ratio = image.width / image.height
    frame_ratio = width / height
    if image_ratio > frame_ratio:
        display_width = width
        display_height = width / image_ratio
        display_left = left
        display_top = top + (height - display_height) / 2
    else:
        display_height = height
        display_width = height * image_ratio
        display_top = top
        display_left = left + (width - display_width) / 2

    slide.shapes.add_picture(str(image_path), display_left, display_top, width=display_width, height=display_height)


def add_content_slide(prs: Presentation, item: dict[str, Any], annotated_path: Path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(248, 250, 252)

    role_chip = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.45),
        Inches(0.35),
        Inches(1.2),
        Inches(0.35),
    )
    role_chip.fill.solid()
    role_chip.fill.fore_color.rgb = RGBColor(37, 99, 235)
    role_chip.line.color.rgb = RGBColor(37, 99, 235)
    add_textbox(slide, Inches(0.62), Inches(0.39), Inches(0.9), Inches(0.18), item["role"], 14, True, (255, 255, 255))

    add_textbox(slide, Inches(0.45), Inches(0.82), Inches(7.6), Inches(0.48), item["title"], 24, True, (15, 23, 42))
    add_textbox(slide, Inches(0.45), Inches(1.25), Inches(7.8), Inches(0.38), item["subtitle"], 14, False, (71, 85, 105))

    image_card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.42),
        Inches(1.72),
        Inches(8.15),
        Inches(5.25),
    )
    image_card.fill.solid()
    image_card.fill.fore_color.rgb = RGBColor(255, 255, 255)
    image_card.line.color.rgb = RGBColor(226, 232, 240)
    fit_image(slide, annotated_path, Inches(0.55), Inches(1.88), Inches(7.88), Inches(4.95))

    note_card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(8.78),
        Inches(1.72),
        Inches(4.05),
        Inches(5.25),
    )
    note_card.fill.solid()
    note_card.fill.fore_color.rgb = RGBColor(255, 255, 255)
    note_card.line.color.rgb = RGBColor(226, 232, 240)
    add_textbox(slide, Inches(9.02), Inches(1.95), Inches(3.2), Inches(0.3), "설명 포인트", 18, True, (30, 41, 59))

    bullet_items = []
    for hotspot in sorted(item.get("hotspots", []), key=lambda row: row["index"]):
        bullet_items.append(f"{hotspot['index']}. {hotspot['label']} - {hotspot['note']}")
    for note in item.get("notes", []):
        bullet_items.append(f"• {note}")
    add_bullets(slide, Inches(9.0), Inches(2.3), Inches(3.5), Inches(4.3), bullet_items[:10])


def build_deck(items: list[dict[str, Any]]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    add_title_slide(prs, len(items))
    for item in items:
      annotated_path = annotate_image(item)
      add_content_slide(prs, item, annotated_path)

    prs.save(str(OUTPUT_PPTX))


def main() -> None:
    ensure_dirs()
    manifests = sorted(MANIFEST_DIR.glob("*.json"))
    if not manifests:
        raise SystemExit("No capture manifest files found. Run the Playwright capture step first.")

    items = [json.loads(path.read_text(encoding="utf-8")) for path in manifests]
    build_deck(items)
    print(OUTPUT_PPTX)


if __name__ == "__main__":
    main()
