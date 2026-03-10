from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches

import generate_training_ppt as base


ROOT = Path(r"C:\Users\baek_\newmso")
OUTPUT_DIR = ROOT / "docs" / "manuals" / "generated"
OUTPUT_FILE = OUTPUT_DIR / "NEWMSO_전체메뉴_상세설명서_2차.pptx"
ANNOTATED_DIR = ROOT / "docs" / "manuals" / "ppt_assets" / "annotated"


MAIN_MENU_SUMMARY = [
    "현재 원본 메인 메뉴: 내정보, 추가기능, 채팅, 게시판, 전자결재, 인사관리, 재고관리, 관리자",
    "조직도는 기능은 존재하지만 현재 고정 메인 메뉴가 아니라 추가기능/검색을 통해 진입",
    "일부 화면은 메인 메뉴 -> 서브메뉴 -> 내부 탭 -> 보조 탭 구조로 들어가야 실제 기능에 도달",
]


DECK_SECTIONS = [
    {
        "title": "메뉴 읽는 법",
        "bullets": [
            "메인 메뉴: 사이드바에서 고르는 가장 바깥 단계",
            "서브메뉴: 사이드바 플라이아웃에 뜨는 두 번째 단계",
            "내부 탭: 화면 안에서 다시 나뉘는 기능 단위",
            "보조 탭: 급여, 근태, 관리자처럼 내부에서 한 번 더 갈라지는 단계",
        ],
    },
    {
        "title": "내정보 상세 구조",
        "image": ANNOTATED_DIR / "03_mypage.png",
        "bullets": [
            "탭: 프로필, 출퇴근, 할 일, 증명서, 급여, 서류제출, 알림",
            "프로필: 역할별 대시보드, 즐겨찾기, 개인 정보 수정",
            "출퇴근: 출근/퇴근 기록 확인, 정정 요청 흐름 연계",
            "알림: 시스템 알림 목록과 알림 설정 확인",
        ],
    },
    {
        "title": "조직도와 추가기능",
        "bullets": [
            "조직도는 현재 고정 메인 메뉴가 아니라 추가기능의 조직도 카드로 들어가는 구조",
            "추가기능 카드: 조직도, 부서별 재고, 근무현황, 인계노트, 퇴원심사, 마감보고, 직원평가",
            "외부 링크: KM Park, U+ 웹팩스",
            "인계노트, 마감보고, 직원평가는 부서/역할/관리자 여부에 따라 제한될 수 있음",
        ],
    },
    {
        "title": "채팅 상세 구조",
        "image": ANNOTATED_DIR / "04_chat.png",
        "bullets": [
            "고정 서브메뉴는 없고 채팅방 목록, 대화 영역, 입력 영역 중심으로 동작",
            "기능: 공지메시지, 1:1, 그룹채팅, 멘션, 파일첨부, 반응, 북마크, 공지등록, 스레드, 투표",
            "슬래시 기능: 연차 요청 초안, 발주 요청 초안 생성",
            "공지메시지 방 작성은 현재 특정 직책에만 허용",
        ],
    },
    {
        "title": "게시판 상세 구조",
        "image": ANNOTATED_DIR / "05_board.png",
        "bullets": [
            "사이드바 플라이아웃 기준: 공지사항, 자유게시판, 경조사, 수술일정, MRI일정",
            "실제 내부 게시판은 익명소리함, 직원제안함까지 포함해 총 7개",
            "수술일정과 MRI일정은 일정형 입력 필드와 템플릿 기능 포함",
            "익명소리함은 현재 관리자/인사 권한 중심으로 조회",
        ],
    },
    {
        "title": "전자결재 전체 구조",
        "image": ANNOTATED_DIR / "06_approval.png",
        "bullets": [
            "서브메뉴: 기안함, 결재함, 작성하기, 캘린더, 양식빌더, 서명관리, 직인관리",
            "작성하기 문서 유형: 인사명령, 연차/휴가, 연차계획서, 연장근무, 물품신청, 수리요청서, 업무기안, 업무협조, 양식신청, 출결정정",
            "기능: 임시저장, 마지막 기안 불러오기, 결재선 템플릿, 일괄 승인/반려",
            "현재 원본 UI에서는 관리성 탭도 함께 보일 수 있음",
        ],
    },
    {
        "title": "인사관리 개요",
        "image": ANNOTATED_DIR / "07_hr.png",
        "bullets": [
            "상위 업무 공간: 인력관리, 근태 · 급여, 복지 · 문서",
            "인사관리는 현재 전체 시스템에서 가장 메뉴가 많은 영역",
            "권한에 따라 보이는 탭이 달라지고, 일부 화면은 보조 탭이 한 번 더 있음",
        ],
    },
    {
        "title": "인사관리 - 인력관리",
        "bullets": [
            "기능 탭: 구성원현황, 인사발령, 포상/징계, 교육, 조직도편집기, 스킬매트릭스, 오프보딩",
            "구성원현황: 직원 목록과 상세 정보 확인",
            "조직도편집기: 조직 구조 수정",
            "오프보딩: 퇴사 절차 관리",
        ],
    },
    {
        "title": "인사관리 - 근태 · 급여",
        "bullets": [
            "기능 탭: 근태, 교대근무, 근무표자동편성, 연차/휴가, 급여, 간호근무표, 공휴일달력",
            "근태 보조 탭: 근태 현황, 연차소멸알림, 지각조퇴분석, 근무형태변경이력, 조기퇴근감지",
            "급여 보조 탭: 급여 메인, 원천징수파일, 4대보험/EDI",
            "근태와 급여는 승인 데이터와 실제 운영 데이터가 연결되는 영역",
        ],
    },
    {
        "title": "인사관리 - 복지 · 문서",
        "bullets": [
            "기능 탭: 건강검진, 경조사, 면허/자격증, 의료기기점검, 비품대여, 사고보고서, 계약, 문서보관함, 증명서, 서류제출, 캘린더",
            "계약 보조 탭: 계약 현황, 계약서 자동생성",
            "캘린더 화면은 공유캘린더와 캘린더동기화 패널이 함께 있음",
            "복지와 문서 기능은 부서 운영과 행정 처리 성격이 강함",
        ],
    },
    {
        "title": "재고관리 전체 구조",
        "image": ANNOTATED_DIR / "08_inventory.png",
        "bullets": [
            "서브메뉴 17개가 현재 사이드바 플라이아웃에 직접 나열됨",
            "현황, 이력, UDI, 명세서, 발주, 스캔, 등록, 자산QR, AS/반품, 거래처, 재고실사, 유통기한 알림, 재고이관, 카테고리, 소모품 통계, 납품 확인서, 수요 예측",
            "등록 화면 내부 방식: form, excel, auto_extract",
            "저재고와 유통기한 임박 품목은 알림과 구매 결재 흐름까지 연결될 수 있음",
        ],
    },
    {
        "title": "관리자 전체 구조",
        "image": ANNOTATED_DIR / "09_admin.png",
        "bullets": [
            "사이드바 플라이아웃 기준: 경영분석, 예산관리, 회사관리, 직원권한, 연차수동부여, 알림자동화, 수술검사템플릿, 팝업관리, 양식빌더, 문서서식, 엑셀등록, 급여이상치, 데이터백업, 데이터초기화, 공문서대장, 감사센터",
            "관리자 화면 내부 탭은 다시 경영분석/감사센터/직접 실행 탭으로 나뉨",
            "현재 원본 코드상 관리자 화면은 사실상 MSO 권한 계정 중심으로 동작",
        ],
    },
    {
        "title": "관리자 - 내부 분석/감사 탭",
        "bullets": [
            "경영분석 내부 탭: 경영대시보드, 재무대시보드, 예산관리, 통합보고서, 법인손익",
            "감사센터 내부 탭: 접근감사로그, 감사로그",
            "직접 실행 탭: 엑셀등록, 알림자동화, 연차수동부여, 회사관리, 직원권한, 수술검사템플릿, 팝업관리, 데이터백업, 데이터초기화, 양식빌더, 문서서식, 급여이상치, 공문서대장",
            "특히 권한, 백업, 초기화, 양식 변경은 운영 영향이 큰 영역",
        ],
    },
    {
        "title": "직원 권한 화면 설명",
        "image": ANNOTATED_DIR / "12_admin_permissions.png",
        "bullets": [
            "직원 목록에서 권한 부여 대상을 먼저 고름",
            "메인 메뉴 권한, 인사 세부 권한, 관리자 권한을 조정",
            "권한 변경 후 실제 메뉴 노출 상태를 반드시 확인하는 것이 안전",
        ],
    },
    {
        "title": "교육과 배포에 쓰는 방법",
        "bullets": [
            "사원 교육: 내정보, 채팅, 게시판, 전자결재 중심",
            "부서장 교육: 전자결재, 인사관리, 재고관리까지 확장",
            "관리자 교육: 관리자 메뉴만 따로 보지 말고 인사/재고/결재와 연결해서 설명",
            "마스터 교육: 기능 설명보다 통제 기준과 변경 승인 절차를 함께 설명",
        ],
    },
]


def add_cover(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(245, 247, 250)

    tag = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.55),
        Inches(0.55),
        Inches(2.4),
        Inches(0.42),
    )
    tag.fill.solid()
    tag.fill.fore_color.rgb = RGBColor(37, 99, 235)
    tag.line.color.rgb = RGBColor(37, 99, 235)
    base.add_textbox(slide, Inches(0.82), Inches(0.61), Inches(2.0), Inches(0.2), "2차 상세 버전", 15, True, (255, 255, 255))

    base.add_textbox(slide, Inches(0.65), Inches(1.35), Inches(7.9), Inches(0.8), "NEWMSO 전체 메뉴 상세 설명서", 28, True, (15, 23, 42))
    base.add_textbox(slide, Inches(0.65), Inches(2.08), Inches(7.6), Inches(0.6), "모메뉴, 서브메뉴, 내부 탭, 보조 탭까지 현재 원본 구조를 기준으로 정리한 PPT입니다.", 16, False, (71, 85, 105))

    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.65),
        Inches(3.0),
        Inches(5.8),
        Inches(2.9),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(255, 255, 255)
    panel.line.color.rgb = RGBColor(226, 232, 240)
    base.add_textbox(slide, Inches(0.95), Inches(3.3), Inches(2.4), Inches(0.3), "이번 버전의 목적", 19, True, (30, 41, 59))
    base.add_bullets(slide, Inches(0.95), Inches(3.72), Inches(5.0), Inches(1.8), MAIN_MENU_SUMMARY)

    right = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(6.8),
        Inches(1.2),
        Inches(5.7),
        Inches(5.7),
    )
    right.fill.solid()
    right.fill.fore_color.rgb = RGBColor(230, 238, 255)
    right.line.color.rgb = RGBColor(191, 219, 254)
    base.add_textbox(slide, Inches(7.1), Inches(1.55), Inches(3.8), Inches(0.35), "포함 범위", 20, True, (29, 78, 216))
    base.add_bullets(
        slide,
        Inches(7.1),
        Inches(2.0),
        Inches(4.9),
        Inches(3.5),
        [
            "내정보",
            "조직도 / 추가기능",
            "채팅",
            "게시판",
            "전자결재",
            "인사관리",
            "재고관리",
            "관리자",
        ],
    )


def add_section_slide(prs: Presentation, section: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(248, 250, 252)

    base.add_textbox(slide, Inches(0.6), Inches(0.6), Inches(7.8), Inches(0.5), section["title"], 24, True, (15, 23, 42))

    if section.get("image"):
        image_card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(0.55),
            Inches(1.25),
            Inches(7.65),
            Inches(5.6),
        )
        image_card.fill.solid()
        image_card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        image_card.line.color.rgb = RGBColor(226, 232, 240)
        base.fit_image(slide, section["image"], Inches(0.75), Inches(1.45), Inches(7.25), Inches(5.2))

        note_card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(8.45),
            Inches(1.25),
            Inches(4.35),
            Inches(5.6),
        )
        note_card.fill.solid()
        note_card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        note_card.line.color.rgb = RGBColor(226, 232, 240)
        base.add_textbox(slide, Inches(8.75), Inches(1.55), Inches(2.6), Inches(0.3), "상세 설명", 18, True, (30, 41, 59))
        base.add_bullets(slide, Inches(8.72), Inches(1.95), Inches(3.65), Inches(4.6), section["bullets"])
    else:
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(0.75),
            Inches(1.45),
            Inches(11.8),
            Inches(4.9),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = RGBColor(226, 232, 240)
        base.add_bullets(slide, Inches(1.05), Inches(1.82), Inches(11.0), Inches(4.2), section["bullets"])


def main():
    base.ensure_dirs()
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(base.SLIDE_W)
    prs.slide_height = Inches(base.SLIDE_H)

    add_cover(prs)
    for section in DECK_SECTIONS:
        add_section_slide(prs, section)

    prs.save(str(OUTPUT_FILE))
    print(OUTPUT_FILE)


if __name__ == "__main__":
    main()
